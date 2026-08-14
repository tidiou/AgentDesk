from pathlib import Path
import pdfplumber
from docx import Document
from pptx import Presentation

from app.schemas.parsed import ParsedDocument


def parse_pdf(filepath: Path) -> ParsedDocument:
    text_parts = []
    with pdfplumber.open(filepath) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

    full_text = "\n".join(text_parts)
    if not full_text.strip():
        raise ValueError(
            f"No extractable text found in '{filepath.name}' — "
            "it may be a scanned/image-based PDF (OCR not supported)"
        )

    return ParsedDocument(
        filename=filepath.name,
        file_type="pdf",
        text=full_text,
        page_count=page_count,
        word_count=len(full_text.split()),
    )


def parse_docx(filepath: Path) -> ParsedDocument:
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)

    # Headings are paragraphs whose style name starts with "Heading"
    sections = [
        p.text for p in doc.paragraphs
        if p.style.name.startswith("Heading") and p.text.strip()
    ]

    return ParsedDocument(
        filename=filepath.name,
        file_type="docx",
        text=full_text,
        sections=sections,
        word_count=len(full_text.split()),
    )


def parse_pptx(filepath: Path) -> ParsedDocument:
    prs = Presentation(filepath)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)

    full_text = "\n".join(t for t in text_parts if t.strip())
    slide_count = len(prs.slides)  # python-pptx supports len() directly on the slide collection

    return ParsedDocument(
        filename=filepath.name,
        file_type="pptx",
        text=full_text,
        page_count=slide_count,
        word_count=len(full_text.split()),
    )


def parse_txt(filepath: Path) -> ParsedDocument:
    full_text = filepath.read_text(encoding="utf-8", errors="replace")
    return ParsedDocument(
        filename=filepath.name,
        file_type="txt",
        text=full_text,
        word_count=len(full_text.split()),
    )


def parse_document(filepath: Path) -> ParsedDocument:
    """Dispatches to the correct parser based on file extension."""
    ext = filepath.suffix.lower()
    dispatch = {
        ".pdf": parse_pdf,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".txt": parse_txt,
    }
    if ext not in dispatch:
        raise ValueError(f"'{ext}' is not a supported document type")
    return dispatch[ext](filepath)